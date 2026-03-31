#!/usr/bin/env python3
"""Build Ninai demo PPTX marketing deck.

Run from backend root:
    python notebooks/build_pptx.py

Output:
    notebooks/ninai_demo.pptx  — 12 slides, 5 case studies, 2 charts
"""
from __future__ import annotations

import io
from pathlib import Path

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import numpy as np

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ── Colour palette ────────────────────────────────────────────────────────────
NAVY      = RGBColor(0x0F, 0x25, 0x47)
TEAL      = RGBColor(0x00, 0xB4, 0xD8)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
DARK      = RGBColor(0x1F, 0x29, 0x37)
MID       = RGBColor(0x6B, 0x72, 0x80)
GREEN_D   = RGBColor(0x14, 0x53, 0x2D)
GREEN_L   = RGBColor(0xDC, 0xFC, 0xE7)
RED_D     = RGBColor(0x7F, 0x1D, 0x1D)
RED_L     = RGBColor(0xFE, 0xE2, 0xE2)
ORANGE    = RGBColor(0xC2, 0x55, 0x00)
YELLOW_L  = RGBColor(0xFF, 0xF3, 0xCD)
GRAY_L    = RGBColor(0xF3, 0xF4, 0xF6)
BLUE_MID  = RGBColor(0x22, 0x3A, 0x5E)

# Slide dimensions (16:9 widescreen)
W = Inches(13.33)
H = Inches(7.5)

# ── Layout helper ─────────────────────────────────────────────────────────────

def _blank_layout(prs: Presentation):
    for layout in prs.slide_layouts:
        if layout.name.lower() in ("blank", ""):
            return layout
    return prs.slide_layouts[6]


def _new_slide(prs: Presentation):
    return prs.slides.add_slide(_blank_layout(prs))


# ── Low-level shape helpers ───────────────────────────────────────────────────

def _rect(slide, left, top, width, height, fill=WHITE, border=None):
    """Add a filled rectangle (autoshape type 1 = MSO rectangle)."""
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if border:
        shape.line.color.rgb = border
        shape.line.width = Pt(0.75)
    else:
        shape.line.color.rgb = fill  # invisible border
    return shape


def _text(slide, text: str, left, top, width, height, *,
          size=14, bold=False, italic=False, color=DARK,
          align=PP_ALIGN.LEFT, v_anchor=MSO_ANCHOR.MIDDLE):
    """Add a text box, handling \\n as paragraph breaks."""
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf = txb.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    tf.vertical_anchor = v_anchor

    lines = str(text).split("\n")
    para = tf.paragraphs[0]
    para.alignment = align
    _run(para, lines[0], size=size, bold=bold, italic=italic, color=color)
    for line in lines[1:]:
        p = tf.add_paragraph()
        p.alignment = align
        _run(p, line, size=size, bold=bold, italic=italic, color=color)
    return txb


def _run(para, text: str, *, size=14, bold=False, italic=False, color=DARK):
    run = para.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return run


# ── Table cell helper ─────────────────────────────────────────────────────────

def _cell(cell, text: str, *,
          bg=None, fg=DARK, size=11, bold=False, align=PP_ALIGN.LEFT):
    if bg:
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    tf = cell.text_frame
    tf.word_wrap = True

    lines = str(text).split("\n")
    para = tf.paragraphs[0]
    para.alignment = align
    _run(para, lines[0], size=size, bold=bold, color=fg)
    for line in lines[1:]:
        p = tf.add_paragraph()
        p.alignment = align
        _run(p, line, size=size, bold=bold, color=fg)


# ── Branded header bar ────────────────────────────────────────────────────────

def _header(slide, title: str, subtitle: str = ""):
    _rect(slide, 0, 0, W, Inches(1.1), fill=NAVY)
    _text(slide, "ninai", Inches(0.3), Inches(0.1), Inches(1.4), Inches(0.55),
          size=22, bold=True, color=TEAL, align=PP_ALIGN.LEFT)
    _text(slide, title, Inches(1.75), Inches(0.08), Inches(9.5), Inches(0.52),
          size=22, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        _text(slide, subtitle, Inches(1.75), Inches(0.62), Inches(10.5), Inches(0.38),
              size=11, color=TEAL, align=PP_ALIGN.LEFT)


# ── Chart generators ──────────────────────────────────────────────────────────

def _chart_anomaly() -> io.BytesIO:
    """Case 4: auth ticket bar chart with baseline."""
    days   = ["Mon", "Tue", "Wed", "Thu", "Fri", "Sat", "Mon\n(now)"]
    counts = [10, 9, 14, 11, 13, 12, 47]
    colors = ["#4A90D9"] * 6 + ["#E53E3E"]

    fig, ax = plt.subplots(figsize=(7.5, 4.0), dpi=130)
    fig.patch.set_facecolor("white")
    bars = ax.bar(days, counts, color=colors, width=0.6, zorder=3, edgecolor="none")
    ax.axhline(12.0, color="#0F2547", linewidth=1.8, linestyle="--",
               label="30-day baseline  (12 tickets/day)", zorder=4)
    ax.bar_label(bars, fontsize=10, fontweight="bold", padding=2)
    ax.annotate(
        "+292% above\nbaseline",
        xy=(6, 47), xytext=(4.6, 54),
        fontsize=9, color="#E53E3E", fontweight="bold",
        ha="center",
        arrowprops=dict(arrowstyle="->", color="#E53E3E", lw=1.5),
    )
    ax.set_ylabel("Auth Tickets", fontsize=11)
    ax.set_title(
        "Case 4 — Auth Service Ticket Spike  (Meridian Technologies)",
        fontsize=12, fontweight="bold", pad=10,
    )
    ax.legend(fontsize=10, framealpha=0.9)
    ax.set_ylim(0, 70)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.grid(axis="y", alpha=0.3, zorder=0)
    plt.tight_layout()

    buf = io.BytesIO()
    plt.savefig(buf, format="png", bbox_inches="tight", dpi=150, facecolor="white")
    plt.close(fig)
    buf.seek(0)
    return buf


def _chart_trend() -> io.BytesIO:
    """Case 5: storage/OS ticket trend with projection and SLA line."""
    days_actual     = [1, 2, 3, 4, 5]
    counts_actual   = [2, 4, 7, 12, 19]
    days_projected  = [5, 6, 7]
    counts_projected = [19, 38, 89]

    fig, ax = plt.subplots(figsize=(7.5, 4.0), dpi=130)
    fig.patch.set_facecolor("white")

    ax.plot(days_actual, counts_actual, "o-", color="#4A90D9", linewidth=2.5,
            markersize=8, label="Observed (Days 1–5)", zorder=5)
    ax.plot(days_projected, counts_projected, "o--", color="#F97316", linewidth=2.5,
            markersize=8, label="Projected (Days 6–7)", zorder=5)
    ax.axhline(50, color="#E53E3E", linewidth=2.0, linestyle="-.",
               label="SLA threshold  (50 tickets/day)", zorder=4)

    ax.annotate(
        "Ninai warns here\n(Day 3 — 4 days early)",
        xy=(3, 7), xytext=(1.1, 25),
        fontsize=9, color="#0F2547", fontweight="bold",
        arrowprops=dict(arrowstyle="->", color="#0F2547", lw=1.4),
        bbox=dict(boxstyle="round,pad=0.3", fc="#E0F2FE", ec="#4A90D9", alpha=0.9),
    )
    ax.annotate(
        "SLA breach — Day 7\n(89 tickets vs limit 50)",
        xy=(7, 89), xytext=(5.4, 72),
        fontsize=9, color="#E53E3E", fontweight="bold",
        arrowprops=dict(arrowstyle="->", color="#E53E3E", lw=1.4),
    )

    ax.set_xlabel("Day", fontsize=11)
    ax.set_ylabel("Storage/OS Tickets", fontsize=11)
    ax.set_title(
        "Case 5 — Exponential Trend Detection  (4-Day Early Warning)",
        fontsize=12, fontweight="bold", pad=10,
    )
    ax.legend(fontsize=10, loc="upper left", framealpha=0.9)
    ax.set_xlim(0.5, 7.5)
    ax.set_ylim(0, 105)
    ax.set_xticks([1, 2, 3, 4, 5, 6, 7])
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.grid(alpha=0.3, zorder=0)
    plt.tight_layout()

    buf = io.BytesIO()
    plt.savefig(buf, format="png", bbox_inches="tight", dpi=150, facecolor="white")
    plt.close(fig)
    buf.seek(0)
    return buf


# ── Individual slides ─────────────────────────────────────────────────────────

def slide_title(prs: Presentation):
    slide = _new_slide(prs)

    # Full navy background
    _rect(slide, 0, 0, W, H, fill=NAVY)
    # Teal accent strip at bottom
    _rect(slide, 0, Inches(6.85), W, Inches(0.65), fill=TEAL)

    # Product name + taglines (left)
    _text(slide, "ninai", Inches(0.55), Inches(0.9), Inches(6.5), Inches(1.3),
          size=72, bold=True, color=TEAL, align=PP_ALIGN.LEFT)
    _text(slide,
          "Intelligent Memory for\nSupport & Engineering Teams",
          Inches(0.55), Inches(2.3), Inches(8.5), Inches(1.1),
          size=24, color=WHITE, align=PP_ALIGN.LEFT)
    _text(slide,
          "5 real scenarios. One system that connects the dots.",
          Inches(0.55), Inches(3.5), Inches(8.5), Inches(0.55),
          size=16, italic=True, color=TEAL, align=PP_ALIGN.LEFT)

    # Right side — case preview cards
    case_labels = [
        ("Case 1", "The Repeat Caller",    "LOW",            GREEN_L, GREEN_D),
        ("Case 2", "The Team Storm",       "LOW–MEDIUM",     GREEN_L, GREEN_D),
        ("Case 3", "The Fix That Wasn't",  "MEDIUM",         YELLOW_L, ORANGE),
        ("Case 4", "The Monday Crash",     "HIGH",           RED_L,  RED_D),
        ("Case 5", "The Slow Boil",        "VERY HIGH",      RED_L,  RED_D),
    ]
    for i, (num, name, complexity, card_bg, card_fg) in enumerate(case_labels):
        y = Inches(0.9 + i * 1.05)
        _rect(slide, Inches(9.5), y, Inches(3.5), Inches(0.85), fill=BLUE_MID)
        _text(slide, num, Inches(9.65), y + Inches(0.04), Inches(0.7), Inches(0.38),
              size=10, bold=True, color=TEAL)
        _text(slide, name, Inches(10.35), y + Inches(0.04), Inches(2.5), Inches(0.38),
              size=12, bold=True, color=WHITE)
        _text(slide, complexity, Inches(9.65), y + Inches(0.44), Inches(3.1), Inches(0.32),
              size=9, italic=True, color=card_bg)

    _text(slide,
          "All demos run on real Kaggle data  ·  Self-hosted  ·  No data leaves your infra",
          Inches(0.55), Inches(6.88), Inches(12.0), Inches(0.45),
          size=12, bold=True, color=NAVY, align=PP_ALIGN.LEFT)


def slide_problem(prs: Presentation):
    slide = _new_slide(prs)
    _header(slide, "The Problem",
            "Why support teams keep fighting the same fires every week")
    _rect(slide, 0, Inches(1.1), W, H - Inches(1.1), fill=GRAY_L)

    pain_points = [
        (
            "Tickets don't talk to each other",
            "A customer calls 3 times about the same broken login. Each CSR sees only their own "
            "ticket. No one connects the pattern. No escalation. No resolution.",
        ),
        (
            "Incidents get closed before they're actually fixed",
            "An incident is marked 'Resolved' at 09:00. By end of day 30 customers still can't "
            "log in. There is no link between the closed ticket and the open complaints.",
        ),
        (
            "Trends go unnoticed until they become a crisis",
            "A ticket subcategory grows  2 → 4 → 7 → 12 → 19  over 5 days. "
            "No single day looks alarming. Day 7 projects to 89 — above SLA limit of 50. "
            "The signal was there. Nobody was looking at the right data.",
        ),
    ]

    for i, (title, body) in enumerate(pain_points):
        y = Inches(1.55 + i * 1.85)
        _rect(slide, Inches(0.38), y, Inches(0.55), Inches(0.55), fill=NAVY)
        _text(slide, str(i + 1), Inches(0.38), y - Inches(0.02), Inches(0.55), Inches(0.55),
              size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        _text(slide, title, Inches(1.1), y - Inches(0.04), Inches(11.8), Inches(0.5),
              size=16, bold=True, color=NAVY)
        _text(slide, body, Inches(1.1), y + Inches(0.48), Inches(11.8), Inches(1.1),
              size=12, color=DARK)


def slide_how_it_works(prs: Presentation):
    slide = _new_slide(prs)
    _header(slide, "How Ninai Works",
            "Memory-native intelligence — each ticket runs a full agent pipeline before a human reads it")
    _rect(slide, 0, Inches(1.1), W, H - Inches(1.1), fill=WHITE)

    # Pipeline boxes
    pipeline = [
        ("Memory\nIngestion",   "Tickets, incidents,\nnotes → structured\nmemory records",        NAVY),
        ("Agent\nPipeline",     "17 specialist agents\nrun on each memory;\nparallel enrichment",  TEAL),
        ("Cross-Record\nFusion","Entity, timeline, \ncredibility & causal\nchain resolved",        RGBColor(0x0E, 0x7C, 0xB4)),
        ("Briefing\nCard",      "CSR / SRE gets the\nfull picture before\nthey press Assign",     GREEN_D),
    ]
    for i, (label, desc, col) in enumerate(pipeline):
        x = Inches(0.5 + i * 3.2)
        _rect(slide, x, Inches(2.0), Inches(2.85), Inches(1.75), fill=col)
        _text(slide, label, x, Inches(2.05), Inches(2.85), Inches(0.7),
              size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        _text(slide, desc, x, Inches(2.77), Inches(2.85), Inches(0.9),
              size=10, color=WHITE, align=PP_ALIGN.CENTER)
        if i < 3:
            _text(slide, "→", x + Inches(2.85), Inches(2.6), Inches(0.35), Inches(0.55),
                  size=22, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    # Agent list
    _text(slide, "17 Specialist Agents:",
          Inches(0.5), Inches(4.0), Inches(4.0), Inches(0.35),
          size=13, bold=True, color=NAVY)
    agent_lines = [
        "EntityResolution  ·  TemporalReasoning  ·  Credibility  ·  FeedbackIntegration",
        "EpisodicGrouping  ·  SiloPropagation  ·  OrgAttention  ·  ConflictDetection",
        "CausalReasoning  ·  AdaptiveConflictResolution  ·  AnomalyDetection",
        "GoalDecomposition  ·  Playbook  ·  PredictiveMonitor  ·  AutonomousGoalGeneration",
        "MetaCognitivePlanning  ·  NarrativeSynthesis",
    ]
    for i, line in enumerate(agent_lines):
        _text(slide, line, Inches(0.5), Inches(4.4 + i * 0.38), Inches(12.5), Inches(0.36),
              size=10, color=MID)

    # Stats bar
    _rect(slide, 0, Inches(6.7), W, Inches(0.8), fill=NAVY)
    stats = [
        ("17",           "specialist agents"),
        ("< 2 s",        "pipeline latency"),
        ("5",            "complexity tiers"),
        ("SLA-aware",    "predictions"),
        ("100%",         "self-hosted"),
    ]
    for i, (val, lbl) in enumerate(stats):
        x = Inches(0.6 + i * 2.6)
        _text(slide, val, x, Inches(6.66), Inches(1.6), Inches(0.42),
              size=15, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
        _text(slide, lbl, x, Inches(7.1), Inches(1.6), Inches(0.32),
              size=9, color=WHITE, align=PP_ALIGN.CENTER)


# ── Reusable comparison-table slide ──────────────────────────────────────────

def _comparison_slide(prs, case_num, title, complexity, scenario,
                      questions, raw_answers, ninai_answers):
    slide = _new_slide(prs)
    _header(slide,
            f"Case {case_num}: {title}",
            f"Complexity: {complexity}  ·  {scenario}")

    n = len(questions)
    table = slide.shapes.add_table(
        n + 1, 3,
        Inches(0.32), Inches(1.22),
        Inches(12.68), Inches(5.88),
    ).table

    table.columns[0].width = Inches(3.5)
    table.columns[1].width = Inches(4.55)
    table.columns[2].width = Inches(4.63)

    _cell(table.cell(0, 0), "Question",       bg=NAVY,                  fg=WHITE, size=12, bold=True, align=PP_ALIGN.CENTER)
    _cell(table.cell(0, 1), "Without Ninai",  bg=RGBColor(0x7F,0x1D,0x1D), fg=WHITE, size=12, bold=True, align=PP_ALIGN.CENTER)
    _cell(table.cell(0, 2), "With Ninai  ✓",  bg=GREEN_D,               fg=WHITE, size=12, bold=True, align=PP_ALIGN.CENTER)

    for i, (q, raw, ninai) in enumerate(zip(questions, raw_answers, ninai_answers), start=1):
        row_bg = RGBColor(0xF8, 0xFA, 0xFC) if i % 2 == 0 else WHITE
        _cell(table.cell(i, 0), q,     bg=row_bg, fg=DARK,  size=11)
        _cell(table.cell(i, 1), raw,   bg=RED_L,  fg=RED_D, size=11)
        _cell(table.cell(i, 2), ninai, bg=GREEN_L, fg=GREEN_D, size=11)

    return slide


def slide_case1(prs):
    return _comparison_slide(
        prs, 1, "The Repeat Caller", "LOW",
        "Emma Chen has called support 3× in 8 days about the same broken login",
        questions=[
            "Is this a new issue?",
            "Has this person called before?",
            "Did the previous fix work?",
            "How urgent is this?",
            "What should the CSR do?",
        ],
        raw_answers=[
            "Looks new — only this ticket is in view",
            "No cross-ticket visibility in a standard portal",
            "Unknown — no follow-up linked to this ticket",
            "Medium — that's the priority label on the form",
            "Start from scratch; attempt another basic reset",
        ],
        ninai_answers=[
            "No — this is contact #3 on the same unresolved login failure",
            "Yes — 3 calls in 8 days; all same root cause",
            "No — password reset (day 1) and cache clear (day 3) both failed within 3 days",
            "CAUTION — repeat-failure pattern triggers escalation protocol",
            "Escalate to L2 immediately. Do not attempt another band-aid fix.",
        ],
    )


def slide_case2(prs):
    return _comparison_slide(
        prs, 2, "The Team Storm", "LOW–MEDIUM",
        "5 DevOps engineers file 8 separate tickets in one week — all converging on one root cause",
        questions=[
            "Are these 8 tickets related?",
            "Is the DevOps team in crisis?",
            "What is the common cause?",
            "Who needs to be alerted urgently?",
            "Should support or engineering own this?",
        ],
        raw_answers=[
            "Unknown — each ticket looks isolated",
            "Unknown — no team-level visibility in the portal",
            "Unknown — each ticket describes a different symptom",
            "Unknown — no cross-team escalation is triggered",
            "Support handles each one; no signal to escalate",
        ],
        ninai_answers=[
            "Yes — all 8 grouped into episode DEVOPS-AUTH-CRISIS-2026",
            "Yes — 5 of 5 DevOps members affected in 7 days",
            "Auth/cert expiry cascading: pipeline → cert → permissions → SSO",
            "VP Engineering + DevOps lead paged; status page updated",
            "Engineering owns this; support closes with a P1 handoff",
        ],
    )


def slide_case3(prs):
    return _comparison_slide(
        prs, 3, "The Fix That Wasn't", "MEDIUM",
        "Incident closed Monday. 4 customer tickets show it is still broken Thursday.",
        questions=[
            "Was the incident really resolved?",
            "Why are customers still reporting failures?",
            "What is the actual root cause?",
            "Who closed the incident incorrectly?",
            "Who needs to be contacted right now?",
        ],
        raw_answers=[
            "Yes — ticket INC-2026-0318 is marked Closed",
            "Unknown — these look like new, separate tickets",
            "Unknown — each ticket just says 'login fails'",
            "Unknown — the DevOps team closed the ticket",
            "Unknown — no escalation target visible from the tickets",
        ],
        ninai_answers=[
            "No — Ninai detected 2 contradictions between close note and live complaints",
            "Key rotation applied to primary node only — replica us-east-1b still uses old key",
            "LDAP users on replica → token validation failure → login loop",
            "Alex Rivera (DevOps) — rotation script skipped replica node",
            "James Kim (security), Alex Rivera (DevOps), Oliver Smith (L2) — all 3 paged",
        ],
    )


def slide_case4_chart(prs):
    slide = _new_slide(prs)
    _header(slide,
            "Case 4: The Monday Morning Crash",
            "Complexity: HIGH  ·  47 tickets in 2 hours. A human is still reading ticket #3.")

    buf = _chart_anomaly()
    slide.shapes.add_picture(buf, Inches(0.4), Inches(1.22), Inches(7.6), Inches(4.55))

    # Info panel (right)
    _rect(slide, Inches(8.3), Inches(1.22), Inches(4.7), Inches(5.5), fill=GRAY_L)
    _text(slide, "What Ninai did in < 2 minutes",
          Inches(8.5), Inches(1.3), Inches(4.4), Inches(0.5),
          size=14, bold=True, color=NAVY)

    facts = [
        ("Detected +292% anomaly vs 30-day baseline",    NAVY),
        ("Identified company-wide scope — 3 departments", NAVY),
        ("Matched playbook MAJOR-INCIDENT-RESPONSE-v3",   NAVY),
        ("Decomposed into 5 subtasks; found blocker",     NAVY),
        ("Declared P1, paged infra on-call at 09:01",     RED_D),
    ]
    icons = ["🔴", "🏢", "📋", "📊", "🚨"]
    for i, ((text, col), icon) in enumerate(zip(facts, icons)):
        y = Inches(1.95 + i * 0.84)
        _text(slide, icon, Inches(8.5), y, Inches(0.45), Inches(0.5), size=16)
        _text(slide, text, Inches(9.0), y + Inches(0.04), Inches(3.8), Inches(0.52),
              size=12, color=col)

    # Pull quote
    _rect(slide, Inches(0.4), Inches(6.05), Inches(7.6), Inches(0.95), fill=NAVY)
    _text(slide,
          '"47 individual complaints are one incident.\nNinai declared it before a human read ticket #3."',
          Inches(0.6), Inches(6.1), Inches(7.2), Inches(0.85),
          size=11, italic=True, color=TEAL, v_anchor=MSO_ANCHOR.MIDDLE)


def slide_case4_table(prs):
    return _comparison_slide(
        prs, 4, "Monday Morning Crash — Response Detail", "HIGH",
        "Auth service authentication failure — 3 departments — Monday 09:00",
        questions=[
            "What is this queue of 47 tickets?",
            "Is this normal volume for a Monday?",
            "Which departments are affected?",
            "What should the on-call do first?",
            "Is there a playbook for this situation?",
        ],
        raw_answers=[
            "A busy Monday — read each one individually",
            "Unknown — no historical baseline visible in the portal",
            "Unknown — each ticket is from one person",
            "Unknown — start reading and manually triage",
            "No automatic matching — on-call must recall from memory",
        ],
        ninai_answers=[
            "+292% spike vs 12/day baseline (30-day avg) — this is a P1 incident",
            "No — 47 tickets in 2 hours vs 12/day average = +292% above normal",
            "Engineering, Operations, Finance — all 3 departments, company-wide",
            "Declare P1 and page infra on-call (blocking subtask #1 of 5)",
            "MAJOR-INCIDENT-RESPONSE-v3 matched at 91% confidence",
        ],
    )


def slide_case5_chart(prs):
    slide = _new_slide(prs)
    _header(slide,
            "Case 5: The Slow Boil",
            "Complexity: VERY HIGH / PREDICTIVE  ·  Detected on Day 3. Crisis would hit Day 7.")

    buf = _chart_trend()
    slide.shapes.add_picture(buf, Inches(0.4), Inches(1.22), Inches(7.6), Inches(4.7))

    # Info panel (right)
    _rect(slide, Inches(8.3), Inches(1.22), Inches(4.7), Inches(5.5), fill=GRAY_L)
    _text(slide, "Ninai's Proactive Warning",
          Inches(8.5), Inches(1.3), Inches(4.4), Inches(0.5),
          size=14, bold=True, color=NAVY)

    stats = [
        ("Day 3",    "Ninai detects exponential trend"),
        ("Day 7",    "Projected breach: 89 tickets (SLA = 50)"),
        ("4 days",   "Lead time before the crisis hits"),
        ("0",        "Human escalations that triggered this"),
        ("78%",      "Prediction confidence"),
    ]
    for i, (val, lbl) in enumerate(stats):
        y = Inches(1.95 + i * 0.84)
        _rect(slide, Inches(8.5), y, Inches(1.1), Inches(0.55), fill=NAVY)
        _text(slide, val, Inches(8.5), y - Inches(0.02), Inches(1.1), Inches(0.58),
              size=14, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
        _text(slide, lbl, Inches(9.7), y + Inches(0.06), Inches(3.1), Inches(0.44),
              size=12, color=DARK)

    _rect(slide, Inches(0.4), Inches(6.1), Inches(7.6), Inches(0.9), fill=NAVY)
    _text(slide,
          '"Ninai warned you on Day 3.\nYou had 4 days to fix it before it became a crisis."',
          Inches(0.6), Inches(6.15), Inches(7.2), Inches(0.8),
          size=11, italic=True, color=TEAL, v_anchor=MSO_ANCHOR.MIDDLE)


def slide_case5_table(prs):
    return _comparison_slide(
        prs, 5, "The Slow Boil — Prediction Accuracy", "VERY HIGH / PREDICTIVE",
        "Storage/OS ticket trend: 2 → 4 → 7 → 12 → 19 over 5 days",
        questions=[
            "Is this trend concerning?",
            "When will it become a crisis?",
            "Was any alert generated?",
            "What should the team do today?",
            "What is the root cause hypothesis?",
        ],
        raw_answers=[
            "No — no single day crosses an alert threshold",
            "Unknown — no forecasting in the ticket portal",
            "No — no SLA threshold crossed yet on any individual day",
            "Continue processing the normal support queue",
            "Unknown — tickets list many different symptoms",
        ],
        ninai_answers=[
            "Yes — exponential growth at 1.9× velocity per 1.5 days",
            "Day 7: projected 89 tickets — SLA breach in 4 days",
            "Yes — AutonomousGoalGenerationAgent raised proactive investigation",
            "Investigate storage/OS health + check this week's OS patches",
            "Storage subsystem + recent OS patches — 78% confidence",
        ],
    )


def slide_summary(prs: Presentation):
    slide = _new_slide(prs)
    _header(slide, "All 5 Cases at a Glance",
            "From a repeat caller to a company-wide crash — one intelligence layer handles them all")

    table = slide.shapes.add_table(
        6, 5,
        Inches(0.3), Inches(1.22),
        Inches(12.73), Inches(5.88),
    ).table

    col_widths = [Inches(2.0), Inches(1.7), Inches(2.5), Inches(3.6), Inches(2.93)]
    for i, w in enumerate(col_widths):
        table.columns[i].width = w

    headers = ["Case", "Complexity", "Key Agents", "What Ninai Detects", "Outcome"]
    for j, h in enumerate(headers):
        _cell(table.cell(0, j), h, bg=NAVY, fg=WHITE, size=12, bold=True, align=PP_ALIGN.CENTER)

    rows = [
        [
            "Case 1\nThe Repeat Caller",
            "LOW",
            "EntityResolution\nTemporal\nCredibility\nFeedback",
            "3rd contact on same unresolved issue — prior fixes failed; escalation required",
            "Escalated to L2 before\nCSR opens the form",
        ],
        [
            "Case 2\nThe Team Storm",
            "LOW–MEDIUM",
            "EpisodicGrouping\nSiloPropagation\nOrgAttention",
            "8 isolated DevOps tickets grouped into 1 auth/cert cascade episode (5 engineers)",
            "VP Engineering paged;\nP1 incident opened",
        ],
        [
            "Case 3\nThe Fix That Wasn't",
            "MEDIUM",
            "ConflictDetect\nCredibility\nCausalReasoning\nAdaptiveConflict",
            "2 contradictions between 'Closed' incident and 4 live customer complaints",
            "Incident reopened;\n3 teams simultaneously paged",
        ],
        [
            "Case 4\nMonday Crash",
            "HIGH",
            "AnomalyDetect\nTemporal\nOrgAttention\nGoalDecomp\nPlaybook",
            "+292% auth ticket spike; company-wide scope; P1 playbook matched at 91%",
            "P1 declared at 09:01 AM\n(ticket #1 of 47)",
        ],
        [
            "Case 5\nSlow Boil",
            "VERY HIGH\nPREDICTIVE",
            "Temporal\nPredictiveMonitor\nAutonomousGoal\nMetaCognitive",
            "Exponential trend on Day 3; SLA breach projected Day 7 (89 vs limit 50)",
            "4-day early warning;\nproactive goal raised",
        ],
    ]

    complexity_fg = [GREEN_D, GREEN_D, ORANGE, RED_D, RED_D]

    for i, (row, cfg) in enumerate(zip(rows, complexity_fg), start=1):
        row_bg = RGBColor(0xF8, 0xFA, 0xFC) if i % 2 == 0 else WHITE
        for j, val in enumerate(row):
            fg = cfg if j == 1 else DARK
            _cell(table.cell(i, j), val, bg=row_bg, fg=fg,
                  size=10, bold=(j == 1))


def slide_cta(prs: Presentation):
    slide = _new_slide(prs)

    _rect(slide, 0, 0, W, H, fill=NAVY)
    _rect(slide, 0, Inches(5.85), W, Inches(1.65), fill=TEAL)

    _text(slide, "ninai", Inches(0.55), Inches(0.6), Inches(5.5), Inches(1.1),
          size=64, bold=True, color=TEAL, align=PP_ALIGN.LEFT)
    _text(slide, "Memory-native intelligence for ops teams.",
          Inches(0.55), Inches(1.85), Inches(9.5), Inches(0.65),
          size=22, color=WHITE, align=PP_ALIGN.LEFT)

    bullets = [
        "Self-hosted — your data never leaves your infrastructure",
        "17 specialist agents — every ticket runs a full intelligence pipeline",
        "SLA-aware predictions — warns you before the threshold is crossed",
        "Connects tickets, incidents and internal notes across all teams",
        "REST API — integrates with any ticket system in hours",
    ]
    for i, b in enumerate(bullets):
        _text(slide, f"✓  {b}",
              Inches(0.75), Inches(2.65 + i * 0.5), Inches(10.5), Inches(0.44),
              size=14, color=WHITE)

    _text(slide,
          "All 5 demo cases run end-to-end in under 10 seconds on local hardware  ·  Open source core",
          Inches(0.55), Inches(5.92), Inches(12.2), Inches(0.48),
          size=13, bold=True, color=NAVY, align=PP_ALIGN.LEFT)
    _text(slide,
          "Try the live notebook demo  ·  Ask for a walkthrough of any case  ·  apache-2.0 licensed",
          Inches(0.55), Inches(6.42), Inches(12.2), Inches(0.42),
          size=12, italic=True, color=NAVY, align=PP_ALIGN.LEFT)


# ── Assembly ──────────────────────────────────────────────────────────────────

def main():
    out = Path(__file__).with_name("ninai_demo.pptx")

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    slide_title(prs)          # 1
    slide_problem(prs)        # 2
    slide_how_it_works(prs)   # 3
    slide_case1(prs)          # 4  Case 1 comparison table
    slide_case2(prs)          # 5  Case 2 comparison table
    slide_case3(prs)          # 6  Case 3 comparison table
    slide_case4_chart(prs)    # 7  Case 4 chart + facts panel
    slide_case4_table(prs)    # 8  Case 4 comparison table
    slide_case5_chart(prs)    # 9  Case 5 chart + stats panel
    slide_case5_table(prs)    # 10 Case 5 comparison table
    slide_summary(prs)        # 11 All-cases summary table
    slide_cta(prs)            # 12 CTA / close

    prs.save(str(out))
    print(f"Wrote {len(prs.slides)} slides -> {out}")


if __name__ == "__main__":
    main()
